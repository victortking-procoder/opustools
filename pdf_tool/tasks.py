# --- tasks.py ---
import os
import logging
import datetime
import shutil
import zipfile
import uuid
import subprocess
from celery import shared_task
from django.conf import settings
from PIL import Image
import json

# Import PDF processing libraries
from pdf2docx import Converter
from pdf2image import convert_from_path
from pptx import Presentation
import tabula
import pandas as pd
from pypdf import PdfWriter, PdfReader # NEW: Import PdfWriter and PdfReader for merging and splitting

# Import the models from the current app
from .models import PdfToolJob, PdfUploadedFile # NEW: Import PdfUploadedFile as well

logger = logging.getLogger(__name__)

@shared_task(bind=True)
def process_file_task(self, job_id):
    """
    Celery task to handle all operations for the pdf_tool based on the 'tool_type'.
    This task consolidates PDF conversion, compression, merging, and splitting.
    """
    job = None
    try:
        job = PdfToolJob.objects.get(id=job_id)
        job.status = 'PROCESSING'
        job.save()

        tool_type = job.tool_type
        logger.info(f"Task {self.request.id}: Starting processing for job ID: {job_id} with tool_type: {tool_type}")

        # Retrieve the list of file paths from the ManyToManyField
        file_path_relative_to_media_root = [
            f.file.name for f in job.uploaded_files.all()
        ]

        # Define the output directory once
        output_dir_relative_to_media_root = os.path.join('pdf_tool_processed', str(job.id))
        os.makedirs(os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root), exist_ok=True)
        
        # Determine if we need to handle a single file or multiple files
        is_single_file_job = tool_type in ['pdf_converter', 'pdf_splitter']
        source_files_abs_paths = [os.path.join(settings.MEDIA_ROOT, path) for path in file_path_relative_to_media_root]

        if tool_type == 'pdf_converter':
            # --- PDF Conversion Logic ---
            if not source_files_abs_paths:
                raise ValueError("No file provided for PDF conversion.")
            
            # For conversion, we assume a single file
            absolute_source_file_path = source_files_abs_paths[0]
            target_format = job.target_format
            
            if not os.path.exists(absolute_source_file_path):
                raise FileNotFoundError(f"Source file does not exist at {absolute_source_file_path}")

            # Define Output File Path
            base_filename = os.path.splitext(os.path.basename(absolute_source_file_path))[0]
            output_extension = target_format.lower()
            output_filename = f"converted_{base_filename}.{output_extension}"
            output_full_absolute_path = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, output_filename)

            # Conversion Logic based on target_format
            if target_format.upper() == 'DOCX':
                cv = Converter(absolute_source_file_path)
                cv.convert(output_full_absolute_path)
                cv.close()
            
            elif target_format.upper() == 'XLSX':
                tables = tabula.read_pdf(absolute_source_file_path, pages='all', multiple_tables=True, pandas_options={'header': None})
                if tables:
                    with pd.ExcelWriter(output_full_absolute_path) as writer:
                        for i, df in enumerate(tables):
                            df = df.dropna(axis=0, how='all').dropna(axis=1, how='all')
                            df.to_excel(writer, sheet_name=f'Table {i+1}', index=False)
                else:
                    df = pd.DataFrame([["No tables found in the PDF."]])
                    with pd.ExcelWriter(output_full_absolute_path) as writer:
                        df.to_excel(writer, sheet_name='Sheet1', index=False, header=False)

            elif target_format.upper() == 'PPTX':
                prs = Presentation()
                images = convert_from_path(absolute_source_file_path)
                for img in images:
                    blank_slide_layout = prs.slide_layouts[6]
                    slide = prs.slides.add_slide(blank_slide_layout)
                    
                    temp_img_path = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, f"temp_{uuid.uuid4()}.jpg")
                    img.save(temp_img_path, 'JPEG')
                    
                    slide.shapes.add_picture(temp_img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
                    os.remove(temp_img_path)
                
                prs.save(output_full_absolute_path)

            elif target_format.upper() == 'JPG':
                images = convert_from_path(absolute_source_file_path)
                image_output_folder = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, 'images')
                os.makedirs(image_output_folder, exist_ok=True)
                
                for i, page in enumerate(images):
                    image_filename = f"page_{i+1}.jpg"
                    image_path = os.path.join(image_output_folder, image_filename)
                    page.save(image_path, 'JPEG')

                zip_filename = f"converted_{base_filename}.zip"
                zip_full_absolute_path = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, zip_filename)
                shutil.make_archive(
                    os.path.splitext(zip_full_absolute_path)[0],
                    'zip',
                    image_output_folder
                )
                shutil.rmtree(image_output_folder)

                download_url = os.path.join(settings.MEDIA_URL, output_dir_relative_to_media_root, zip_filename).replace('\\', '/')
                job.output_url = download_url
                job.status = 'COMPLETED'
                job.save()
                return

            else:
                raise ValueError(f"Unsupported target format: {target_format}")

            # Update Job with Output URL and Status (for single-file outputs)
            download_url = os.path.join(settings.MEDIA_URL, output_dir_relative_to_media_root, output_filename).replace('\\', '/')
            job.output_url = download_url
            job.status = 'COMPLETED'
            job.save()
            return # IMPORTANT: Add return here to prevent fall-through

        elif tool_type == 'file_compressor':
            # --- PDF Compression Logic using Ghostscript ---
            if not source_files_abs_paths:
                raise ValueError("No files provided for PDF compression.")

            compression_level = job.compression_level
            if not compression_level:
                raise ValueError("No compression level specified for the job.")

            if compression_level == 'high':
                gs_options = ['-dPDFSETTINGS=/screen']
            elif compression_level == 'medium':
                gs_options = ['-dPDFSETTINGS=/ebook']
            else: # low
                gs_options = ['-dPDFSETTINGS=/printer']

            compressed_file_paths = []
            
            for absolute_source_file_path in source_files_abs_paths:
                if not absolute_source_file_path.lower().endswith('.pdf'):
                    raise ValueError(f"File {absolute_source_file_path} is not a PDF.")

                if not os.path.exists(absolute_source_file_path):
                    raise FileNotFoundError(f"Source file does not exist: {absolute_source_file_path}")
                
                base_filename = os.path.splitext(os.path.basename(absolute_source_file_path))[0]
                output_filename = f"compressed_{base_filename}.pdf"
                output_full_absolute_path = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, output_filename)
                
                gs_command = [
                    'gs',
                    '-sDEVICE=pdfwrite',
                    '-dCompatibilityLevel=1.4',
                    '-dNOPAUSE',
                    '-dBATCH',
                    '-dSAFER',
                    '-sOutputFile=' + output_full_absolute_path,
                    absolute_source_file_path
                ]
                gs_command[3:3] = gs_options

                try:
                    subprocess.run(gs_command, check=True, capture_output=True, text=True)
                except subprocess.CalledProcessError as e:
                    logger.error(f"Ghostscript command failed: {e.stderr}")
                    raise RuntimeError(f"PDF compression failed: {e.stderr}")

                compressed_file_paths.append(output_full_absolute_path)

            if len(compressed_file_paths) > 1:
                zip_filename = "compressed_files.zip"
                zip_full_absolute_path = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, zip_filename)
                
                with zipfile.ZipFile(zip_full_absolute_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                    for file_path in compressed_file_paths:
                        zipf.write(file_path, os.path.basename(file_path))
                
                for file_path in compressed_file_paths:
                    os.remove(file_path)

                download_url = os.path.join(settings.MEDIA_URL, output_dir_relative_to_media_root, zip_filename).replace('\\', '/')
            else:
                output_file_path = compressed_file_paths[0]
                output_file_name = os.path.basename(output_file_path)
                download_url = os.path.join(settings.MEDIA_URL, output_dir_relative_to_media_root, output_file_name).replace('\\', '/')

            job.output_url = download_url
            job.status = 'COMPLETED'
            job.save()
            return

        # --- NEW: PDF Merging Logic ---
        elif tool_type == 'pdf_merger':
            if not source_files_abs_paths or len(source_files_abs_paths) < 2:
                raise ValueError("Merging requires at least two PDF files.")

            # Get the merge order (stored as JSON string in job object)
            merge_order = job.merge_order
            if not merge_order:
                raise ValueError("Merge order is required for PDF merging jobs.")

            try:
                merge_order = json.loads(merge_order)  # ✅ parse string into list
            except json.JSONDecodeError as e:
                job.status = 'FAILED'
                job.error_message = f"Invalid merge_order format: {e}"
                job.save()
                return

            # Map original filenames to absolute paths
            file_path_map = {os.path.basename(path): path for path in source_files_abs_paths}

            pdf_merger = PdfWriter()
            output_filename = "merged_document.pdf"
            output_full_absolute_path = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, output_filename)

            try:
                # Add PDFs in specified order
                for filename in merge_order:
                    if filename in file_path_map:
                        pdf_merger.append(file_path_map[filename])
                    else:
                        logger.warning(f"File {filename} not found in the uploaded files. Skipping.")

                # Write out merged PDF
                with open(output_full_absolute_path, 'wb') as output_file:
                    pdf_merger.write(output_file)
            finally:
                pdf_merger.close()

            download_url = os.path.join(settings.MEDIA_URL, output_dir_relative_to_media_root, output_filename).replace('\\', '/')
            job.output_url = download_url
            job.status = 'COMPLETED'
            job.save()
            return

            
        # --- NEW: PDF Splitting Logic ---
        elif tool_type == 'pdf_splitter':
            if not source_files_abs_paths or len(source_files_abs_paths) != 1:
                job.status = 'FAILED'
                job.error_message = "Splitting requires exactly one PDF file."
                job.save()
                return

            absolute_source_file_path = source_files_abs_paths[0]
            if not os.path.exists(absolute_source_file_path):
                job.status = 'FAILED'
                job.error_message = f"Source file does not exist at {absolute_source_file_path}"
                job.save()
                return

            page_ranges_str = job.page_ranges
            if not page_ranges_str:
                job.status = 'FAILED'
                job.error_message = "Page ranges are required for PDF splitting jobs."
                job.save()
                return

            split_pdf_paths = []
            base_filename = os.path.splitext(os.path.basename(absolute_source_file_path))[0]

            try:
                with open(absolute_source_file_path, 'rb') as source_file:
                    pdf_reader = PdfReader(source_file)
                    num_pages = len(pdf_reader.pages)

                    # --- Parse and validate page ranges ---
                    parsed_ranges = []
                    try:
                        for part in page_ranges_str.split(','):
                            part = part.strip()
                            if not part:
                                continue  # skip empty tokens like ",,"

                            if '-' in part and part.count('-') == 1:
                                start_str, end_str = part.split('-')
                                if not start_str.strip() or not end_str.strip():
                                    raise ValueError(f"Invalid range format: '{part}'")

                                start, end = int(start_str), int(end_str)
                                if start > end:
                                    raise ValueError(f"Start page {start} cannot be greater than end page {end}")
                                if start < 1 or end > num_pages:
                                    raise ValueError(f"Page range {start}-{end} is out of bounds. PDF has {num_pages} pages.")
                                parsed_ranges.append((start, end))
                            else:
                                # Treat as a single page
                                page_number = int(part)
                                if page_number < 1 or page_number > num_pages:
                                    raise ValueError(f"Page {page_number} is out of bounds. PDF has {num_pages} pages.")
                                parsed_ranges.append((page_number, page_number))
                    except (ValueError, IndexError) as e:
                        job.status = 'FAILED'
                        job.error_message = f"Invalid page range format: {e}"
                        job.save()
                        return

                    # --- Process each range ---
                    for start, end in parsed_ranges:
                        pdf_writer = PdfWriter()
                        # Add pages (user input is 1-based, PdfReader is 0-based)
                        for page_num in range(start - 1, end):
                            pdf_writer.add_page(pdf_reader.pages[page_num])

                        # Output filename
                        if start == end:
                            output_filename = f"{base_filename}_page_{start}.pdf"
                        else:
                            output_filename = f"{base_filename}_pages_{start}-{end}.pdf"

                        output_full_absolute_path = os.path.join(
                            settings.MEDIA_ROOT, output_dir_relative_to_media_root, output_filename
                        )
                        with open(output_full_absolute_path, 'wb') as output_file:
                            pdf_writer.write(output_file)

                        split_pdf_paths.append(output_full_absolute_path)

            except Exception as e:
                job.status = 'FAILED'
                job.error_message = f"Processing failed due to an internal error: {e}"
                job.save()
                return

            # --- Zip all the split files together ---
            zip_filename = f"split_{base_filename}.zip"
            zip_full_absolute_path = os.path.join(settings.MEDIA_ROOT, output_dir_relative_to_media_root, zip_filename)
            with zipfile.ZipFile(zip_full_absolute_path, 'w', zipfile.ZIP_DEFLATED) as zipf:
                for file_path in split_pdf_paths:
                    zipf.write(file_path, os.path.basename(file_path))

            # Clean up individual split files
            for file_path in split_pdf_paths:
                os.remove(file_path)

            download_url = os.path.join(settings.MEDIA_URL, output_dir_relative_to_media_root, zip_filename).replace('\\', '/')
            job.output_url = download_url
            job.status = 'COMPLETED'
            job.save()
            return

        
    except PdfToolJob.DoesNotExist:
        logger.error(f"Task {self.request.id}: PdfToolJob with ID {job_id} does not exist.", exc_info=True)
    except FileNotFoundError as e:
        logger.error(f"Task {self.request.id}: FileNotFoundError for job {job_id}: {e}", exc_info=True)
        if job:
            job.status = 'FAILED'
            job.error_message = f"Processing failed: Source file not found. Error: {e}"
            job.save()
    except Exception as e:
        logger.error(f"Task {self.request.id}: An unexpected error occurred for job {job_id}: {e}", exc_info=True)
        if job:
            job.status = 'FAILED'
            job.error_message = f"Processing failed due to an internal error: {str(e)}"
            job.save()